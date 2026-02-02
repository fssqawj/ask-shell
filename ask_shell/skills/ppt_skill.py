"""PPT Generation Skill - Create PowerPoint presentations"""

import os
import re
import json
from loguru import logger
from datetime import datetime
from typing import List, Optional, Dict, Any, Callable
from .base_skill import BaseSkill, SkillExecutionResponse, SkillCapability
from ..llm.base import BaseLLMClient
from ..llm.openai_client import OpenAIClient

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE = None
    MSO_AUTO_SHAPE_TYPE = None


class PPTSkill(BaseSkill):
    """
    Skill for generating PowerPoint presentations
    
    This skill can:
    1. Create PPT from text outlines
    2. Generate slides with specific themes
    3. Add charts and images to presentations
    """
    
    def __init__(self):
        super().__init__()
        self.initialized = Presentation is not None and MSO_AUTO_SHAPE_TYPE is not None
        self.output_dir = "./output"
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir, exist_ok=True)
        
        # Initialize LLM client for content generation
        try:
            self.llm: BaseLLMClient = OpenAIClient()
            self.system_prompt = """你是一个专业的PPT内容策划师。用户会给你一个主题和任务要求，以及可能的历史交互信息。请为PowerPoint演示文稿生成合适的大纲和每页的详细内容。

你的回复必须是一个JSON对象，格式如下：
{
    "title": "演示文稿的整体标题",
    "outline": [
        {
            "title": "幻灯片标题",
            "content": "幻灯片详细内容，包括要点、说明等",
            "layout_type": "content_type",  // 可选值：'title_only', 'title_content', 'section_header', 'two_content', 'list', 'bullet_points', 'image_placeholder'
            "elements": [  // 可选：定义额外的视觉元素
                {
                    "type": "icon",  // 如 'checkmark', 'warning', 'info', 'arrow', 'star'
                    "position": "top_left",  // 或 'top_right', 'bottom_left', 'bottom_right'
                    "size": "medium"  // 或 'small', 'large'
                }
            ]
        },
        ...
    ]
}

要求：
1. 幻灯片数量通常为3-8张，根据内容复杂程度调整
2. 内容要有层次感，从概述到细节逐步展开
3. 考虑历史信息，使内容连贯一致
4. 确保内容专业、有深度且易于理解
5. 适应用户的具体需求和任务背景
6. **美化要求**：内容应适合美观的PPT展示，包含清晰的标题、要点分明的内容，适合视觉呈现，考虑使用列表、要点、短句等形式便于PPT美化排版
7. **结构化内容**：使用markdown-style formatting (bold **text**, bullet points, etc.) to enhance visual appeal
8. **Layout considerations**：Specify appropriate layout_type for each slide based on content (e.g., 'list' for bullet points, 'bullet_points' for key points, 'image_placeholder' when visual elements would help)"""
            self.llm.set_system_prompt(self.system_prompt)
            self.llm.set_direct_mode(True)  # Use direct mode for content generation
            self.llm_available = True
        except Exception as e:
            # If LLM initialization fails (e.g., missing API key), still allow basic functionality
            print(f"Warning: LLM initialization failed: {str(e)}. PPT skill will use basic functionality.")
            self.llm_available = False
    
    def get_capabilities(self) -> List[SkillCapability]:
        """PPT skill provides file generation capability"""
        return [SkillCapability.FILE_GENERATION]
    
    def execute(
        self,
        task: str,
        context: Optional[Dict[str, Any]] = None,
        stream_callback=None,
        **kwargs
    ) -> SkillExecutionResponse:
        """
        Execute PPT generation
            
        Args:
            task: Task description (e.g., "create a PPT about AI")
            context: Execution context
            **kwargs: Additional parameters including selection_reasoning
            
        Returns:
            SkillExecutionResponse with generated file path
        """
        # Get the reasoning for why this skill was selected
        selection_reasoning = kwargs.get('selection_reasoning', '')
        
        if not self.initialized:
            return SkillExecutionResponse(
                thinking="Checking if required library (python-pptx) is installed",
                direct_response="Error: python-pptx library is not installed or incomplete. Please install it using: pip install python-pptx",
                generated_files=[],
                file_metadata={"status": "missing_dependency", "required_library": "python-pptx"}
            )
        
        try:
            # Get execution context from context
            last_result = context.get('last_result')
            history = context.get('history', [])
            
            # Generate presentation outline using LLM based on task and context
            title, outline = self._generate_outline_with_llm(task, history, last_result, stream_callback)
            
            # Generate the presentation
            filename = self._generate_presentation(title, outline)
            
            # Create detailed direct_response with presentation content
            content_summary = f"\n\nPresentation Content Summary:\nTitle: {title}\nSlides:\n"
            for i, slide_data in enumerate(outline, 1):
                content_summary += f"  {i}. {slide_data['title']}\n     Content: {slide_data['content'][:100]}{'...' if len(slide_data['content']) > 100 else ''}\n"
            
            return SkillExecutionResponse(
                thinking=f"Creating PowerPoint presentation '{title}' with {len(outline)} slides based on the given task and context",
                direct_response=f"Successfully created PowerPoint presentation '{filename}' with {len(outline)} slides.{content_summary}",
                generated_files=[filename],
                file_metadata={
                    "status": "success",
                    "slides_count": len(outline),
                    "title": title,
                    "filename": filename
                }
            )
        except Exception as e:
            return SkillExecutionResponse(
                thinking=f"Failed to create PowerPoint presentation: {str(e)}",
                direct_response=f"Error creating PowerPoint presentation: {str(e)}",
                generated_files=[],
                file_metadata={"status": "error", "error": str(e)}
            )
    
    def _generate_outline_with_llm(self, task: str, history: List[Any], last_result: Optional[Any], stream_callback: Optional[Callable] = None) -> tuple[str, List[Dict[str, Any]]]:
        """Generate presentation outline using LLM based on task and context"""
        
        # If LLM is not available, fall back to basic parsing
        if not self.llm_available:
            return self._parse_task_basic(task, history, last_result)
        
        context_str = f"{task}\n\n"
        
        message = f"{context_str}\n\n请根据以上信息为PowerPoint演示文稿生成合适的大纲和每页的详细内容。" 
        
        try:
            # Call LLM to generate outline with direct parsing using PPTSkillResponse dataclass
            logger.info(f"PPT Skill LLM Prompt: {message}")
            
            from ..models.types import PPTSkillResponse
            # Generate and directly parse into PPTSkillResponse
            llm_response = self.llm.generate(message, last_result, stream_callback=stream_callback, history=history, response_class=PPTSkillResponse)
            logger.info(f"PPT Skill LLM Response: {llm_response}")
            
            # If the response is already parsed (when response_class is provided), use it directly
            if hasattr(llm_response, 'outline'):  # It's already a PPTSkillResponse object
                parsed_response = llm_response
            else:
                # Fallback to raw JSON parsing if needed
                parsed = None
                try:
                    parsed = json.loads(llm_response.raw_json)
                except json.JSONDecodeError:
                    logger.opt(exception=True).error("Error parsing JSON from LLM response")
                    # Fallback to extracting JSON from raw response if direct parsing fails
                    response_text = llm_response.raw_json
                    start_idx = response_text.find('{')
                    end_idx = response_text.rfind('}')
                    
                    if start_idx != -1 and end_idx != -1 and start_idx < end_idx:
                        json_str = response_text[start_idx:end_idx+1]
                        try:
                            parsed = json.loads(json_str)
                        except json.JSONDecodeError:
                            logger.opt(exception=True).error("Error parsing JSON from LLM response")
                            pass
                
                if parsed:
                    # Create PPTSkillResponse manually
                    title = parsed.get("title", "Generated Presentation")
                    outline = parsed.get("outline", [])
                    parsed_response = PPTSkillResponse(title=title, outline=outline)
                else:
                    # If JSON parsing failed, use default values
                    parsed_response = PPTSkillResponse(title="Generated Presentation", outline=[])
            
            title = parsed_response.title or "Generated Presentation"
            outline = parsed_response.outline or []
            
            if title and outline:
                # Validate outline structure
                if not isinstance(outline, list) or len(outline) == 0:
                    # Fallback to a default structure
                    outline = [
                        {"title": title, "content": "Introduction slide for the presentation"},
                        {"title": "Overview", "content": "Main points of the presentation"},
                        {"title": "Details", "content": "Detailed information and analysis"},
                        {"title": "Conclusion", "content": "Summary and next steps"}
                    ]
                
                return title.strip(), outline
            
            # If JSON parsing failed, fallback to default behavior
            return "Generated Presentation", [
                {"title": "Introduction", "content": f"Presentation about: {task}"},
                {"title": "Overview", "content": "Main points of the presentation"},
                {"title": "Details", "content": "Detailed information based on context"},
                {"title": "Conclusion", "content": "Summary and next steps"}
            ]
        except Exception as e:
            logger.opt(exception=e).error("Error generating outline with LLM")
            # Fallback to default behavior if LLM call fails
            return "Generated Presentation", [
                {"title": "Introduction", "content": f"Presentation about: {task}"},
                {"title": "Overview", "content": "Main points of the presentation"},
                {"title": "Details", "content": "Detailed information based on context"},
                {"title": "Conclusion", "content": "Summary and next steps"}
            ]
    
    def _parse_task_basic(self, task: str, history: List[Any], last_result: Optional[Any]) -> tuple[str, List[Dict[str, Any]]]:
        """Basic task parsing when LLM is not available"""
        # Extract title from task
        title = task.split('about')[-1].strip() if 'about' in task.lower() else task.split('presentation')[-1].strip() if 'presentation' in task.lower() else task
        title = title or "Generated Presentation"
        
        # Create basic outline based on the task and any context from history
        outline = [
            {"title": title, "content": f"Introduction to {title}"},
            {"title": "Background", "content": "Context and importance of this topic"},
            {"title": "Key Points", "content": "Main ideas and concepts"},
            {"title": "Applications", "content": "Practical uses and examples"},
            {"title": "Future Outlook", "content": "Trends and developments"},
            {"title": "Conclusion", "content": "Summary and recommendations"}
        ]
        
        # Adjust outline based on history if available
        if history:
            context_summary = " "
            for result in history[-2:]:  # Use last 2 history items
                if hasattr(result, 'skill_response') and hasattr(result.skill_response, 'direct_response'):
                    context_summary += f" {result.skill_response.direct_response}"
                elif hasattr(result, 'command'):
                    context_summary += f" Command: {result.command}"
            
            if len(context_summary) > 2:
                outline[2]["content"] = f"Main ideas and concepts based on context:{context_summary}"
        
        # Limit to 4-6 slides for readability
        return title.strip(), outline[:6]
    
    def _generate_presentation(self, title: str, outline: List[Dict[str, Any]]) -> str:
        """Generate the actual PowerPoint presentation with enhanced styling, layouts, and text overflow handling"""
        prs = Presentation()
            
        # Apply theme and color scheme
        # Use a professional color scheme
        
            
        # Add title slide with enhanced styling
        title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_placeholder = title_slide.shapes.title
        subtitle_placeholder = title_slide.placeholders[1]
            
        # Style the title slide
        if title_placeholder:
            title_placeholder.text = title
            title_frame = title_placeholder.text_frame
            if title_frame.paragraphs:
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
                title_para.alignment = PP_ALIGN.CENTER
            
        if subtitle_placeholder:
            subtitle_placeholder.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}\nBy Ask-Shell PPT Skill"
            subtitle_frame = subtitle_placeholder.text_frame
            if subtitle_frame.paragraphs:
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(18)
                subtitle_para.font.color.rgb = RGBColor(102, 102, 102)  # Gray
                subtitle_para.alignment = PP_ALIGN.CENTER
            
        # Add content slides with varied layouts and styling
        for i, slide_data in enumerate(outline):
            # Determine layout based on slide_data or default to cycling
            layout_type = slide_data.get("layout_type", "title_content")
                        
            # Initialize layout_index for fallback
            layout_index = (i + 1) % 3
                        
            # Map layout types to actual slide layouts
            if layout_type in ["title_content", "content", "bullet_points", "list"]:
                content_slide_layout = prs.slide_layouts[1]  # Title and Content
            elif layout_type == "section_header":
                content_slide_layout = prs.slide_layouts[2]  # Section Header
            elif layout_type == "two_content":
                content_slide_layout = prs.slide_layouts[3]  # Two Content
            else:
                # Default to cycling through layouts
                if layout_index == 0:
                    content_slide_layout = prs.slide_layouts[1]  # Title and Content
                elif layout_index == 1:
                    content_slide_layout = prs.slide_layouts[2]  # Section Header
                else:
                    content_slide_layout = prs.slide_layouts[3]  # Two Content
                
            slide = prs.slides.add_slide(content_slide_layout)
                
            # Style the title
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = slide_data["title"][0:255] if len(slide_data["title"]) > 255 else slide_data["title"]
                title_frame = title_shape.text_frame
                if title_frame.paragraphs:
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(32)
                    title_para.font.bold = True
                    title_para.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
                
            # Process content to handle text overflow and formatting
            content_text = slide_data["content"]
                
            # Add icons or visual elements if specified
            elements = slide_data.get("elements", [])
            for element in elements:
                if element.get("type") == "icon":
                    self._add_icon_to_slide(slide, element)
                
            # Handle content based on layout type
            if layout_type in ["title_content", "content", "bullet_points", "list"]:
                # Layout 1 typically has title (index 0) and content (index 1)
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]  # Usually the second placeholder
                    # Split content into chunks if it's too long
                    formatted_content = self._format_content_for_ppt(content_text)
                    content_placeholder.text = formatted_content
                                
                    # Style the content text
                    content_frame = content_placeholder.text_frame
                    content_frame.word_wrap = True
                    content_frame.fit_text = True  # Auto-fit text to placeholder
                                
                    # Style paragraphs in content
                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black
                        paragraph.line_spacing = 1.3
                        
            elif layout_type == "section_header":
                # Layout 2 typically has title (index 0) and subtitle (index 2)
                if len(slide.placeholders) > 2:
                    content_placeholder = slide.placeholders[2]  # Usually the subtitle placeholder
                    # Split content into chunks if it's too long
                    formatted_content = self._format_content_for_ppt(content_text)
                    content_placeholder.text = formatted_content
                                
                    # Style the content text
                    content_frame = content_placeholder.text_frame
                    content_frame.word_wrap = True
                    content_frame.fit_text = True  # Auto-fit text to placeholder
                                
                    # Style paragraphs in content
                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black
                        paragraph.line_spacing = 1.3
                        
            else:  # Two Content layout (layout 3) or fallback
                # Layout 3 typically has title (index 0), content1 (index 1), and content2 (index 2)
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]  # First content area
                    # Split content into chunks if it's too long
                    formatted_content = self._format_content_for_ppt(content_text)
                    content_placeholder.text = formatted_content
                                
                    # Style the content text
                    content_frame = content_placeholder.text_frame
                    content_frame.word_wrap = True
                    content_frame.fit_text = True  # Auto-fit text to placeholder
                                
                    # Style paragraphs in content
                    for paragraph in content_frame.paragraphs:
                        paragraph.font.size = Pt(18)
                        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black
                        paragraph.line_spacing = 1.3
                            
                # Add secondary content to second placeholder if available
                if len(slide.placeholders) > 2:
                    secondary_placeholder = slide.placeholders[2]
                    secondary_content = "Supporting details or examples related to the main content"
                    if layout_type == "two_content":
                        # If it's specifically a two_content layout, use content from the data if available
                        parts = content_text.split("\n\n")
                        if len(parts) > 1:
                            secondary_content = "\n".join(parts[1:]) if len(parts) > 1 else secondary_content
                    secondary_placeholder.text = secondary_content
                                
                    secondary_frame = secondary_placeholder.text_frame
                    secondary_frame.word_wrap = True
                    secondary_frame.fit_text = True
                                
                    for paragraph in secondary_frame.paragraphs:
                        paragraph.font.size = Pt(16)
                        paragraph.font.color.rgb = RGBColor(50, 50, 50)
                        paragraph.line_spacing = 1.2
            
        # Add a thank you/conclusion slide
        thank_you_layout = prs.slide_layouts[0]
        thank_you_slide = prs.slides.add_slide(thank_you_layout)
        thank_title = thank_you_slide.shapes.title
        thank_subtitle = thank_you_slide.placeholders[1]
            
        if thank_title:
            thank_title.text = "Thank You"
            title_frame = thank_title.text_frame
            if title_frame.paragraphs:
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 51, 102)
                title_para.alignment = PP_ALIGN.CENTER
            
        if thank_subtitle:
            thank_subtitle.text = "Questions & Discussion"
            subtitle_frame = thank_subtitle.text_frame
            if subtitle_frame.paragraphs:
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(24)
                subtitle_para.font.color.rgb = RGBColor(102, 102, 102)
                subtitle_para.alignment = PP_ALIGN.CENTER
            
        # Generate filename
        safe_title = re.sub(r'[\/*?:"<>|]', "", title)[:50]  # Remove invalid filename characters
        filename = os.path.join(self.output_dir, f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
            
        # Save the presentation
        prs.save(filename)
            
        return filename
        
    def _format_content_for_ppt(self, content: str) -> str:
        """Format content for PPT display, handling text overflow and cleaning markdown"""
        # Limit content length to prevent overflow
        max_length = 800  # Reduced from 1000 to prevent overflow
            
        # Clean up markdown-like formatting for PPT display
        cleaned_content = content
            
        # Remove bold markers (**text**) but keep the text
        import re
        cleaned_content = re.sub(r'\*\*(.*?)\*\*', r'\1', cleaned_content)  # Remove **bold**
        cleaned_content = re.sub(r'__(.*?)__', r'\1', cleaned_content)  # Remove __underline__
            
        # Replace markdown-style bullet points with proper bullets
        lines = cleaned_content.split('\n')
        formatted_lines = []
        for line in lines:
            # Handle markdown list items
            if line.strip().startswith('- '):
                formatted_lines.append(line.strip())
            elif line.strip().startswith('* '):
                formatted_lines.append(line.strip())
            elif line.strip().startswith('#'):
                # Skip markdown headers
                continue
            else:
                formatted_lines.append(line)
            
        cleaned_content = '\n'.join(formatted_lines)
            
        # If content is still too long, truncate it and add indicator
        if len(cleaned_content) > max_length:
            cleaned_content = cleaned_content[:max_length].rsplit(' ', 1)[0] + "... (continued on next slide if needed)"
            
        return cleaned_content
        
    def _add_icon_to_slide(self, slide, element):
        """Add an icon or visual element to the slide"""
        try:
            # Check if required enums are available
            if MSO_AUTO_SHAPE_TYPE is None:
                return  # Skip if library not available
                
            icon_type = element.get("type", "")
            position = element.get("position", "top_left")
            size = element.get("size", "medium")
                
            # Define icon mappings
            icon_map = {
                "checkmark": MSO_AUTO_SHAPE_TYPE.CHECKBOX, 
                "warning": MSO_AUTO_SHAPE_TYPE.TRIANGLE,
                "info": MSO_AUTO_SHAPE_TYPE.OVAL,
                "arrow": MSO_AUTO_SHAPE_TYPE.PENTAGON,
                "star": MSO_AUTO_SHAPE_TYPE.STAR_8_POINT
            }
                
            # Determine size
            size_map = {
                "small": (Inches(0.3), Inches(0.3)),
                "medium": (Inches(0.5), Inches(0.5)),
                "large": (Inches(0.8), Inches(0.8))
            }
            width, height = size_map.get(size, size_map["medium"])
                
            # Determine position coordinates
            left_map = {"top_left": Inches(0.2), "top_right": Inches(8.8), "bottom_left": Inches(0.2), "bottom_right": Inches(8.8)}
            top_map = {"top_left": Inches(1.0), "top_right": Inches(1.0), "bottom_left": Inches(5.0), "bottom_right": Inches(5.0)}
                
            left = left_map.get(position, Inches(0.2))
            top = top_map.get(position, Inches(1.0))
                
            # Add shape based on icon type
            if icon_type in icon_map:
                shape_type = icon_map[icon_type]
                shape = slide.shapes.add_shape(shape_type, left, top, width, height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 100, 200)  # Blue color
                shape.line.color.rgb = RGBColor(0, 50, 150)  # Darker blue border
        except Exception as e:
            # If adding icons fails, continue without them
            pass
        
    def get_description(self) -> str:
        """Get skill description"""
        return "专业PPT制作工具，可以根据需求创建PowerPoint演示文稿，支持自动生成标题页和内容页"
